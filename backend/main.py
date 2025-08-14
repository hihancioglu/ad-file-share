import os
import shutil
import secrets
from datetime import datetime, timedelta
import mimetypes
import zipfile

# Ensure previewed file types have proper MIME types
mimetypes.add_type("image/png", ".png")
mimetypes.add_type("image/jpeg", ".jpg")
mimetypes.add_type("image/jpeg", ".jpeg")
mimetypes.add_type("application/vnd.ms-excel", ".xls")
mimetypes.add_type(
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"
)
mimetypes.add_type("application/msword", ".doc")
mimetypes.add_type(
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".docx",
)
mimetypes.add_type("application/vnd.ms-powerpoint", ".ppt")
mimetypes.add_type(
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".pptx",
)
mimetypes.add_type("text/csv", ".csv")

import msal
import requests
from functools import lru_cache

from flask import (
    Flask,
    jsonify,
    render_template,
    request,
    send_file,
    Response,
    session,
    redirect,
)
from flask_cors import CORS
import ldap3
from dotenv import load_dotenv
from database import SessionLocal, add_missing_columns
from models import (
    ShareLink,
    DownloadLog,
    Team,
    TeamMember,
    TeamFile,
    Notification,
    Activity,
    UserShare,
    UserFile,
    FileMessage,
)
from sqlalchemy import func

load_dotenv()
add_missing_columns()

PUBLIC_BASE_URL = os.getenv("PUBLIC_BASE_URL", "https://sendv2.baylan.info.tr").rstrip("/")

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "super-secret-key")
CORS(app, supports_credentials=True)

# Allow uploads up to 1 GB
MAX_UPLOAD_SIZE = 1024 * 1024 * 1024  # 1 GB
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_SIZE

DATA_DIR = "data"
os.makedirs(DATA_DIR, exist_ok=True)


def get_unique_filename(directory: str, filename: str) -> str:
    base, ext = os.path.splitext(filename)
    candidate = filename
    counter = 1
    while os.path.exists(os.path.join(directory, candidate)):
        candidate = f"{base}_{counter}{ext}"
        counter += 1
    return candidate


# Track in-progress uploads to handle name collisions for chunked uploads
current_uploads = {}


def format_file_size(num_bytes: int) -> str:
    size = float(num_bytes)
    for unit in ["B", "KB", "MB", "GB", "TB"]:
        if size < 1000 or unit == "TB":
            if unit == "B":
                return f"{int(size)} B"
            return f"{size:.1f} {unit}"
        size /= 1000


def format_remaining(td: timedelta) -> str:
    total = int(td.total_seconds())
    if total <= 0:
        return "0 saat"
    days, rem = divmod(total, 86400)
    hours = rem // 3600
    if days >= 7:
        return f"{days} gün"
    if days > 0:
        return f"{days} gün {hours} saat"
    return f"{hours} saat"


LDAP_SERVER = os.getenv("LDAP_SERVER")
LDAP_DOMAIN = os.getenv("LDAP_DOMAIN")
LDAP_USER = os.getenv("LDAP_USER")
LDAP_PASSWORD = os.getenv("LDAP_PASSWORD")
LDAP_BASE_DN = os.getenv("LDAP_BASE_DN", "")
LDAP_SEARCH_FILTER = os.getenv(
    "LDAP_SEARCH_FILTER", "(&(objectClass=user)(sAMAccountName=*{query}*))"
)
# Only expose these OUs in the user selection tree
ALLOWED_OUS = {"BAYLAN3", "BAYLAN4", "BAYLAN5"}

ADMIN_USERS = {
    u.strip().lower() for u in os.getenv("ADMIN_USERS", "").split(",") if u.strip()
}

LOGIN_WHITELIST = {
    u.strip().lower() for u in os.getenv("WHITE_LIST", "").split(",") if u.strip()
}


def is_admin(username: str) -> bool:
    return username.lower() in ADMIN_USERS

GRAPH_TENANT_ID = os.getenv("GRAPH_TENANT_ID")
GRAPH_CLIENT_ID = os.getenv("GRAPH_CLIENT_ID")
GRAPH_CLIENT_SECRET = os.getenv("GRAPH_CLIENT_SECRET")
GRAPH_SENDER = os.getenv("GRAPH_SENDER", "")


def get_user_names(username: str):
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{LDAP_USER}"
    try:
        conn = ldap3.Connection(
            server,
            user=user_dn,
            password=LDAP_PASSWORD,
            authentication=ldap3.NTLM,
        )
        if not conn.bind():
            return "", "", ""
        search_filter = f"(&(objectClass=user)(sAMAccountName={username}))"
        conn.search(LDAP_BASE_DN, search_filter, attributes=["givenName", "sn"])
        if conn.entries:
            entry = conn.entries[0]
            given = getattr(entry, "givenName", None)
            sn = getattr(entry, "sn", None)
            return (
                given.value if given else "",
                sn.value if sn else "",
            )
    except Exception:
        pass
    finally:
        try:
            conn.unbind()
        except Exception:
            pass
    return "", ""


def has_manager(username: str) -> bool:
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{LDAP_USER}"
    try:
        conn = ldap3.Connection(
            server,
            user=user_dn,
            password=LDAP_PASSWORD,
            authentication=ldap3.NTLM,
        )
        if not conn.bind():
            return False
        search_filter = f"(&(objectClass=user)(sAMAccountName={username}))"
        conn.search(LDAP_BASE_DN, search_filter, attributes=["manager"])
        if conn.entries:
            entry = conn.entries[0]
            manager = getattr(entry, "manager", None)
            return bool(manager and manager.value)
    except Exception:
        pass
    finally:
        try:
            conn.unbind()
        except Exception:
            pass
    return False


def get_full_name(username: str):
    given, sn = get_user_names(username)
    full_name = f"{given} {sn}".strip()
    return full_name or username


def create_notification(username: str, message: str, team_id=None):
    db = SessionLocal()
    try:
        db.add(Notification(username=username, message=message, team_id=team_id))
        db.commit()
    finally:
        db.close()


def log_activity(usernames, message: str, category: str = "general"):
    if isinstance(usernames, str):
        usernames = [usernames]
    db = SessionLocal()
    try:
        for u in usernames:
            msg = message
            prefix = f"{u} kullanıcısı "
            if msg.startswith(prefix):
                msg = msg[len(prefix):]
            db.add(Activity(username=u, message=msg, category=category))
        db.commit()
    finally:
        db.close()


def get_manager_info(username: str):
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{LDAP_USER}"
    try:
        conn = ldap3.Connection(
            server,
            user=user_dn,
            password=LDAP_PASSWORD,
            authentication=ldap3.NTLM,
        )
        if not conn.bind():
            return "", "", ""
        search_filter = f"(&(objectClass=user)(sAMAccountName={username}))"
        conn.search(LDAP_BASE_DN, search_filter, attributes=["manager"])
        if not conn.entries:
            return "", "", ""
        manager_dn = getattr(conn.entries[0], "manager", None)
        if not manager_dn:
            return "", "", ""
        conn.search(
            manager_dn.value,
            "(objectClass=user)",
            attributes=["mail", "sAMAccountName", "displayName"],
        )
        if conn.entries:
            mail_attr = getattr(conn.entries[0], "mail", None)
            user_attr = getattr(conn.entries[0], "sAMAccountName", None)
            name_attr = getattr(conn.entries[0], "displayName", None)
            return (
                user_attr.value if user_attr else "",
                mail_attr.value if mail_attr else "",
                name_attr.value if name_attr else "",
            )
    except Exception:
        return "", "", ""
    finally:
        try:
            conn.unbind()
        except Exception:
            pass
    return "", "", ""


def get_user_email(username: str) -> str:
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{LDAP_USER}"
    try:
        conn = ldap3.Connection(
            server,
            user=user_dn,
            password=LDAP_PASSWORD,
            authentication=ldap3.NTLM,
        )
        if not conn.bind():
            return ""
        search_filter = f"(&(objectClass=user)(sAMAccountName={username}))"
        conn.search(LDAP_BASE_DN, search_filter, attributes=["mail"])
        if conn.entries:
            mail_attr = getattr(conn.entries[0], "mail", None)
            return mail_attr.value if mail_attr else ""
    except Exception:
        return ""
    finally:
        try:
            conn.unbind()
        except Exception:
            pass
    return ""

def authenticate_user(username: str, password: str) -> bool:
    """Authenticate a user against LDAP."""
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{username}"
    try:
        conn = ldap3.Connection(
            server, user=user_dn, password=password, authentication=ldap3.NTLM
        )
        if not conn.bind():
            return False
        return True
    except Exception:
        return False
    finally:
        try:
            conn.unbind()
        except Exception:
            pass


@app.route("/manager/name", methods=["GET"])
def manager_name_endpoint():
    user = session.get("username")
    if not user:
        return jsonify(manager="")
    _, _, manager_name = get_manager_info(user)
    if not manager_name:
        return jsonify(manager="")
    return jsonify(manager=manager_name)


def require_manager_auth(link):
    """Ensure the requester is the manager of the link's owner or an admin."""
    user = session.get("username")
    if not user:
        return redirect(f"/login?next={request.path}")
    manager_user, _, _ = get_manager_info(link.username)
    if user != manager_user and not is_admin(user):
        return render_template("message.html", message="Yetkisiz"), 403
    return None


def send_approval_email(
    username: str,
    filename: str,
    approve_token: str,
    reject_token: str,
    approver_email: str,
    purpose: str,
):
    manager_email = approver_email
    if not (
        manager_email
        and GRAPH_TENANT_ID
        and GRAPH_CLIENT_ID
        and GRAPH_CLIENT_SECRET
        and GRAPH_SENDER
    ):
        return
    approval_link = f"{PUBLIC_BASE_URL}/share/approve/{approve_token}"
    reject_link = f"{PUBLIC_BASE_URL}/share/reject/{reject_token}"
    subject = "Dosya Paylaşımı Onayı"
    full_name = get_full_name(username)
    body = (
        f"<p>Baylan Send Dosya Paylaşım Platformu</p>"
        f"<p>'{full_name}' kullanıcısı '{filename}' dosyasını herkese açık olarak paylaşmak istiyor.</p>"
        f"<p>Bu bağlantıya sahip olan 3. kişiler dosyayı indirebilir.</p>"
        f"<p>Kullanım amacı: {purpose}</p>"
        f"<p>"
        f"<a href='{approval_link}' style='padding:10px 20px; background-color:#4CAF50; color:white; text-decoration:none;'>Onayla</a>"
        f"<a href='{reject_link}' style='padding:10px 20px; background-color:#f44336; color:white; text-decoration:none; margin-left:10px;'>Reddet</a>"
        f"</p>"
    )
    authority = f"https://login.microsoftonline.com/{GRAPH_TENANT_ID}"
    app = msal.ConfidentialClientApplication(
        GRAPH_CLIENT_ID, authority=authority, client_credential=GRAPH_CLIENT_SECRET
    )
    result = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    token_value = result.get("access_token")
    if not token_value:
        return
    message = {
        "message": {
            "subject": subject,
            "body": {"contentType": "HTML", "content": body},
            "toRecipients": [{"emailAddress": {"address": manager_email}}],
        }
    }
    headers = {
        "Authorization": f"Bearer {token_value}",
        "Content-Type": "application/json",
    }
    try:
        requests.post(
            f"https://graph.microsoft.com/v1.0/users/{GRAPH_SENDER}/sendMail",
            headers=headers,
            json=message,
        )
    except Exception:
        pass


def find_share_token(username: str, filename: str):
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(username=username, filename=filename).first()
        if link and link.expires_at and link.expires_at < datetime.utcnow():
            db.delete(link)
            db.commit()
            return None
        if link and link.rejected:
            return None
        return link.token if link else None
    finally:
        db.close()


def create_share_link(
    token: str,
    username: str,
    filename: str,
    expires_at=None,
    approved: bool = False,
    approve_token: str | None = None,
    reject_token: str | None = None,
    purpose: str = "",
):
    db = SessionLocal()
    try:
        db.add(
            ShareLink(
                token=token,
                approve_token=approve_token,
                reject_token=reject_token,
                username=username,
                filename=filename,
                expires_at=expires_at,
                approved=approved,
                rejected=False,
                purpose=purpose,
            )
        )
        db.commit()
    finally:
        db.close()


def delete_share_link(username: str, filename: str):
    db = SessionLocal()
    try:
        db.query(ShareLink).filter_by(username=username, filename=filename).delete()
        db.commit()
    finally:
        db.close()


def delete_share_notification(username: str, filename: str):
    mgr_user, _, _ = get_manager_info(username)
    if not mgr_user:
        mgr_user = next(iter(ADMIN_USERS), None)
        if not mgr_user:
            return
    db = SessionLocal()
    try:
        message = f"'{filename}' dosyası için onay bekleyen paylaşım"
        db.query(Notification).filter_by(username=mgr_user, message=message).delete()
        db.commit()
    finally:
        db.close()


@lru_cache(maxsize=1024)
def get_country_from_ip(ip_address: str) -> str:
    """Return the country name for a given IP address.

    The primary lookup uses ipapi.co. If that fails to return a country
    (due to network errors, rate limits, etc.), a secondary lookup is
    attempted using ipwho.is. Results are cached to reduce repeated
    lookups for the same IP address.
    """

    # Try ipapi.co first
    try:
        resp = requests.get(
            f"https://ipapi.co/{ip_address}/json/", timeout=5
        )
        if resp.ok:
            data = resp.json()
            country = data.get("country_name")
            if country:
                return country
    except Exception:
        pass

    # Fallback to ipwho.is
    try:
        resp = requests.get(f"https://ipwho.is/{ip_address}", timeout=5)
        if resp.ok:
            data = resp.json()
            country = data.get("country")
            if country:
                return country
    except Exception:
        pass

    return ""


def log_download(username: str, filename: str, downloader: str | None = None):
    ip_addr = request.headers.get("X-Forwarded-For", request.remote_addr)
    country = get_country_from_ip(ip_addr) if ip_addr else ""
    db = SessionLocal()
    try:
        db.add(
            DownloadLog(
                username=username,
                filename=filename,
                ip_address=ip_addr,
                country=country,
            )
        )
        db.commit()
    finally:
        db.close()
    if downloader and downloader != username:
        log_activity(
            username,
            f"{downloader} kullanıcısı '{filename}' dosyanı indirdi",
            "download",
        )


def set_file_expiry(username: str, filename: str, expires_dt, description: str = ""):
    db = SessionLocal()
    try:
        meta = (
            db.query(UserFile)
            .filter_by(username=username, filename=filename)
            .first()
        )
        if meta:
            meta.expires_at = expires_dt
            meta.description = description
        else:
            db.add(
                UserFile(
                    username=username,
                    filename=filename,
                    expires_at=expires_dt,
                    description=description,
                )
            )
        db.commit()
    finally:
        db.close()


def cleanup_expired_files():
    now = datetime.utcnow()
    db = SessionLocal()
    try:
        expired = (
            db.query(UserFile)
            .filter(UserFile.deleted_at != None)
            .filter(UserFile.deleted_at < now - timedelta(days=15))
            .all()
        )
        for meta in expired:
            trash_path = os.path.join(DATA_DIR, "_trash", meta.username, meta.filename)
            if os.path.exists(trash_path):
                os.remove(trash_path)
            db.query(ShareLink).filter_by(
                username=meta.username, filename=meta.filename
            ).delete()
            db.delete(meta)
        db.commit()
    finally:
        db.close()


@app.route("/", methods=["GET"])
def read_app():
    if "username" not in session:
        return redirect("/login")
    return render_template("app.html")


@app.route("/login", methods=["GET"])
def read_login():
    if "username" in session:
        return redirect("/")
    return render_template("login.html")


@app.route("/login", methods=["POST"])
def login():
    username = request.form.get("username", "").lower()
    password = request.form.get("password")
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{username}"
    try:
        conn = ldap3.Connection(
            server, user=user_dn, password=password, authentication=ldap3.NTLM
        )
        if not conn.bind():
            return jsonify(success=False, error="Kullanıcı adı veya şifre hatalı")
        conn.unbind()
        if not has_manager(username) and username not in LOGIN_WHITELIST:
            return jsonify(
                success=False,
                error="Sisteme giriş için Bilgi İşlemi arayınız.",
            )
        session["username"] = username
        given, sn = get_user_names(username)
        log_activity(
            username,
            f"{username} kullanıcısı sisteme giriş yaptı",
            "login",
        )
        return jsonify(
            success=True,
            username=username,
            givenName=given,
            sn=sn,
            admin=is_admin(username),
        )
    except Exception as e:
        return jsonify(success=False, error=str(e))


@app.route("/logout", methods=["POST"])
def logout():
    """Clear the user session and log the logout event."""
    username = session.pop("username", None)
    if username:
        log_activity(
            username,
            f"{username} kullanıcısı sistemden çıktı",
            "logout",
        )
    session.clear()
    return jsonify(success=True)


@app.route("/users/list", methods=["GET"])
def list_users():
    query = request.args.get("q", "")
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{LDAP_USER}"
    try:
        conn = ldap3.Connection(
            server,
            user=user_dn,
            password=LDAP_PASSWORD,
            authentication=ldap3.NTLM,
        )
        if not conn.bind():
            return jsonify(success=False, error="LDAP bağlantısı başarısız")
        search_filter = LDAP_SEARCH_FILTER.format(query=query)
        conn.search(LDAP_BASE_DN, search_filter, attributes=["sAMAccountName"])
        users = [e.sAMAccountName.value for e in conn.entries]
        return jsonify(users=users)
    except Exception as e:
        return jsonify(success=False, error=str(e))
    finally:
        try:
            conn.unbind()
        except Exception:
            pass


def _list_users_in_dn(conn, dn: str):
    conn.search(
        dn,
        "(objectClass=user)",
        search_scope=ldap3.LEVEL,
        attributes=["sAMAccountName", "givenName", "sn"],
    )
    users = []
    for e in conn.entries:
        username = getattr(e, "sAMAccountName", None)
        given = getattr(e, "givenName", None)
        sn = getattr(e, "sn", None)
        users.append(
            {
                "username": username.value if username else "",
                "givenName": given.value if given else "",
                "sn": sn.value if sn else "",
            }
        )
    return users


def _build_ou_tree(conn, base_dn: str):
    conn.search(
        base_dn,
        "(objectClass=organizationalUnit)",
        search_scope=ldap3.LEVEL,
        attributes=["ou"],
    )
    entries = list(conn.entries)
    tree = []
    for entry in entries:
        name = entry.ou.value if "ou" in entry else ""
        if "pc" in name.lower():
            continue
        dn = entry.entry_dn
        node = {
            "name": name,
            "users": _list_users_in_dn(conn, dn),
            "children": _build_ou_tree(conn, dn),
        }
        tree.append(node)
    return tree


@app.route("/users/tree", methods=["GET"])
def users_tree():
    server = ldap3.Server(LDAP_SERVER)
    user_dn = f"{LDAP_DOMAIN}\\{LDAP_USER}"
    try:
        conn = ldap3.Connection(
            server,
            user=user_dn,
            password=LDAP_PASSWORD,
            authentication=ldap3.NTLM,
        )
        if not conn.bind():
            return jsonify(success=False, error="LDAP bağlantısı başarısız")
        tree = _build_ou_tree(conn, LDAP_BASE_DN)
        tree = [node for node in tree if node["name"] in ALLOWED_OUS]
        return jsonify(tree=tree)
    except Exception as e:
        return jsonify(success=False, error=str(e))
    finally:
        try:
            conn.unbind()
        except Exception:
            pass

@app.route("/upload", methods=["POST"])
def upload_file():
    """Upload endpoint supporting chunked uploads.

    Chunks are uploaded with fields:
    - username: owner of the file
    - filename: original file name
    - chunk_index: index of this chunk (0-based)
    - total_chunks: total number of chunks
    - file: binary data for this chunk

    Legacy single-request uploads are still supported for small files.
    """
    cleanup_expired_files()
    username = request.form.get("username")
    expires_at = request.form.get("expires_at")
    expires_dt = datetime.strptime(expires_at, "%Y-%m-%d") if expires_at else None
    user_dir = os.path.join(DATA_DIR, username)
    os.makedirs(user_dir, exist_ok=True)

    # Parameters for chunked upload
    filename = request.form.get("filename")
    chunk = request.files.get("file")
    chunk_index = request.form.get("chunk_index")
    total_chunks = request.form.get("total_chunks")
    description = request.form.get("description", "")

    if filename and chunk and chunk_index is not None and total_chunks:
        chunk_index = int(chunk_index)
        total_chunks = int(total_chunks)
        key = (username, filename)
        if chunk_index == 0:
            final_name = get_unique_filename(user_dir, filename)
            current_uploads[key] = final_name
        final_name = current_uploads.get(key, filename)
        file_path = os.path.join(user_dir, final_name)
        mode = "wb" if chunk_index == 0 else "ab"

        data = chunk.read()
        current_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        if current_size + len(data) > MAX_UPLOAD_SIZE:
            current_uploads.pop(key, None)
            return (
                jsonify(success=False, error="Dosya 1 GB boyut sınırını aşıyor"),
                413,
            )
        with open(file_path, mode) as f:
            f.write(data)
        # If this was the last chunk, respond with completion info
        if chunk_index + 1 == total_chunks:
            current_uploads.pop(key, None)
            set_file_expiry(username, final_name, expires_dt, description)
            log_activity(
                username,
                f"{username} kullanıcısı '{final_name}' dosyasını yükledi",
                "upload",
            )
            return jsonify(success=True, filenames=[final_name])
        return jsonify(success=True, chunk_index=chunk_index)

    # Fallback to legacy upload for already assembled files
    files = request.files.getlist("files")
    if not files:
        file = request.files.get("file")
        files = [file] if file else []

    uploaded = []
    for file in files:
        if file and file.filename:
            if file.content_length and file.content_length > MAX_UPLOAD_SIZE:
                return (
                    jsonify(success=False, error="Dosya 1 GB boyut sınırını aşıyor"),
                    413,
                )
            final_name = get_unique_filename(user_dir, file.filename)
            file_path = os.path.join(user_dir, final_name)
            file.save(file_path)
            uploaded.append(final_name)
            set_file_expiry(username, final_name, expires_dt, description)
            log_activity(
                username,
                f"{username} kullanıcısı '{final_name}' dosyasını yükledi",
                "upload",
            )

    return jsonify(success=True, filenames=uploaded)


@app.route("/list", methods=["POST"])
def list_files():
    cleanup_expired_files()
    username = request.form.get("username")
    admin_mode = request.form.get("admin") and is_admin(username)
    if not admin_mode:
        user_dir = os.path.join(DATA_DIR, username)
        if not os.path.exists(user_dir):
            return jsonify(files=[])

        db = SessionLocal()
        try:
            metas = {
                m.filename: (m.expires_at, m.description or "")
                for m in db.query(UserFile)
                .filter_by(username=username)
                .filter(UserFile.deleted_at == None)
                .all()
            }
            now = datetime.utcnow()
            links = {
                l.filename: {
                    "token": l.token,
                    "expires_at": l.expires_at,
                    "approved": l.approved,
                    "rejected": l.rejected,
                }
                for l in db.query(ShareLink)
                .filter_by(username=username)
                .filter(ShareLink.rejected == False)
                .filter((ShareLink.expires_at == None) | (ShareLink.expires_at > now))
                .all()
            }
            counts = {
                fn: cnt
                for fn, cnt in db.query(DownloadLog.filename, func.count())
                .filter_by(username=username)
                .group_by(DownloadLog.filename)
                .all()
            }
            msg_counts = {
                fn: cnt
                for fn, cnt in db.query(FileMessage.filename, func.count())
                .filter_by(username=username, read=False)
                .group_by(FileMessage.filename)
                .all()
            }
        finally:
            db.close()

        _, _, manager_name = get_manager_info(username)
        if not manager_name:
            admin_user = next(iter(ADMIN_USERS), None)
            manager_name = get_full_name(admin_user) if admin_user else ""
        files = []
        for filename in os.listdir(user_dir):
            file_path = os.path.join(user_dir, filename)
            stat = os.stat(file_path)
            exp, desc = metas.get(filename, (None, ""))
            link_info = links.get(filename, {})
            token = link_info.get("token")
            link_exp = link_info.get("expires_at")
            approved = link_info.get("approved", False)
            rejected = link_info.get("rejected", False)
            mgr_name = manager_name if token and not approved else ""
            files.append(
                {
                    "title": filename,
                    "added": datetime.fromtimestamp(stat.st_mtime).strftime(
                        "%d/%m/%Y %H:%M:%S"
                    ),
                    "extension": os.path.splitext(filename)[1].lstrip("."),
                    "description": desc,
                    "size": stat.st_size,
                    "expires_at": exp.strftime("%d/%m/%Y") if exp else "",
                    "public_expires_at": link_exp.strftime("%d/%m/%Y")
                    if link_exp
                    else "",
                    "expires_in": format_remaining(exp - now) if exp else "",
                    "public_expires_in": format_remaining(link_exp - now)
                    if link_exp
                    else "",
                    "link": f"{PUBLIC_BASE_URL}/public/{token}" if token and not rejected else "",
                    "approved": approved,
                    "rejected": rejected,
                    "manager_name": mgr_name,
                    "download_count": counts.get(filename, 0),
                    "message_count": msg_counts.get(filename, 0),
                }
            )
        files.sort(key=lambda f: f["added"], reverse=True)
        return jsonify(files=files)

    db = SessionLocal()
    try:
        metas = {
            (m.username, m.filename): (m.expires_at, m.description or "")
            for m in db.query(UserFile)
            .filter(UserFile.deleted_at == None)
            .all()
        }
        now = datetime.utcnow()
        links = {
            (l.username, l.filename): {
                "token": l.token,
                "expires_at": l.expires_at,
                "approved": l.approved,
                "rejected": l.rejected,
            }
            for l in db.query(ShareLink)
            .filter(ShareLink.rejected == False)
            .filter((ShareLink.expires_at == None) | (ShareLink.expires_at > now))
            .all()
        }
        counts = {
            (username, filename): cnt
            for username, filename, cnt in db.query(
                DownloadLog.username,
                DownloadLog.filename,
                func.count(),
            )
            .group_by(DownloadLog.username, DownloadLog.filename)
            .all()
        }
        msg_counts = {
            (username, filename): cnt
            for username, filename, cnt in db.query(
                FileMessage.username,
                FileMessage.filename,
                func.count(),
            )
            .filter(FileMessage.read == False)
            .group_by(FileMessage.username, FileMessage.filename)
            .all()
        }
    finally:
        db.close()

    files = []
    for user in os.listdir(DATA_DIR):
        # Skip the special trash directory to avoid showing deleted files
        if user == "_trash":
            continue
        user_dir = os.path.join(DATA_DIR, user)
        if not os.path.isdir(user_dir):
            continue
        _, _, manager_name = get_manager_info(user)
        if not manager_name:
            admin_user = next(iter(ADMIN_USERS), None)
            manager_name = get_full_name(admin_user) if admin_user else ""
        for filename in os.listdir(user_dir):
            file_path = os.path.join(user_dir, filename)
            stat = os.stat(file_path)
            exp, desc = metas.get((user, filename), (None, ""))
            link_info = links.get((user, filename), {})
            token = link_info.get("token")
            link_exp = link_info.get("expires_at")
            approved = link_info.get("approved", False)
            rejected = link_info.get("rejected", False)
            mgr_name = manager_name if token and not approved else ""
            files.append(
                {
                    "title": filename,
                    "username": user,
                    "added": datetime.fromtimestamp(stat.st_mtime).strftime(
                        "%d/%m/%Y %H:%M:%S"
                    ),
                    "extension": os.path.splitext(filename)[1].lstrip("."),
                    "description": desc,
                    "size": stat.st_size,
                    "expires_at": exp.strftime("%d/%m/%Y") if exp else "",
                    "public_expires_at": link_exp.strftime("%d/%m/%Y")
                    if link_exp
                    else "",
                    "expires_in": format_remaining(exp - now) if exp else "",
                    "public_expires_in": format_remaining(link_exp - now)
                    if link_exp
                    else "",
                    "link": f"{PUBLIC_BASE_URL}/public/{token}" if token and not rejected else "",
                    "approved": approved,
                    "rejected": rejected,
                    "manager_name": mgr_name,
                    "download_count": counts.get((user, filename), 0),
                    "message_count": msg_counts.get((user, filename), 0),
                }
            )
    files.sort(key=lambda f: f["added"], reverse=True)
    return jsonify(files=files)


@app.route("/download", methods=["POST"])
def download_file():
    cleanup_expired_files()
    username = request.form.get("username")
    filename = request.form.get("filename")
    user_dir = os.path.join(DATA_DIR, username)
    file_path = os.path.join(user_dir, filename)
    if not os.path.exists(file_path):
        return jsonify(success=False, error="Dosya bulunamadı")
    log_download(username, filename, session.get("username"))
    return send_file(file_path, as_attachment=True, download_name=filename)


@app.route("/public/<token>/message", methods=["POST"])
def public_message(token):
    sender = request.form.get("sender", "")
    text = request.form.get("message", "").strip()
    if not text:
        return jsonify(success=False, error="Mesaj boş olamaz"), 400
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(token=token).first()
        if (
            not link
            or link.rejected
            or not link.approved
            or (link.expires_at and link.expires_at < datetime.utcnow())
        ):
            return jsonify(success=False, error="Bağlantı geçersiz"), 404
        username = link.username
        filename = link.filename
        db.add(
            FileMessage(
                token=token,
                username=username,
                filename=filename,
                sender=sender,
                message=text,
            )
        )
        db.commit()
    finally:
        db.close()
    create_notification(username, f"'{filename}' dosyasına yeni mesaj geldi")
    return jsonify(success=True)


@app.route("/messages/list", methods=["POST"])
def list_messages():
    username = request.form.get("username")
    filename = request.form.get("filename")
    db = SessionLocal()
    try:
        msgs = (
            db.query(FileMessage)
            .filter_by(username=username, filename=filename)
            .order_by(FileMessage.created_at)
            .all()
        )
        result = [
            {
                "id": m.id,
                "sender": m.sender,
                "message": m.message,
                "created_at": m.created_at.strftime("%d/%m/%Y %H:%M"),
            }
            for m in msgs
        ]
        for m in msgs:
            m.read = True
        db.commit()
    finally:
        db.close()
    return jsonify(messages=result)


@app.route("/messages/delete", methods=["POST"])
def delete_message():
    msg_id = request.form.get("id")
    db = SessionLocal()
    try:
        msg = db.query(FileMessage).filter_by(id=msg_id).first()
        if msg:
            db.delete(msg)
            db.commit()
            return jsonify(success=True)
        return jsonify(success=False, error="Mesaj bulunamadı"), 404
    finally:
        db.close()


@app.route("/download/logs", methods=["POST"])
def download_logs():
    username = request.form.get("username")
    filename = request.form.get("filename")
    db = SessionLocal()
    try:
        logs = (
            db.query(DownloadLog)
            .filter_by(username=username, filename=filename)
            .order_by(DownloadLog.timestamp.desc())
            .all()
        )
        data = []
        for log in logs:
            ts = log.timestamp + timedelta(hours=3)
            data.append(
                {
                    "timestamp": ts.strftime("%d/%m/%Y %H:%M:%S"),
                    "ip_address": log.ip_address,
                    "country": log.country,
                }
            )
        return jsonify(logs=data)
    finally:
        db.close()


@app.route("/delete", methods=["POST"])
def delete_file():
    username = request.form.get("username")
    filename = request.form.get("filename")
    user_dir = os.path.join(DATA_DIR, username)
    file_path = os.path.join(user_dir, filename)
    if not os.path.exists(file_path):
        return jsonify(success=False, error="Dosya bulunamadı")
    trash_dir = os.path.join(DATA_DIR, "_trash", username)
    os.makedirs(trash_dir, exist_ok=True)
    shutil.move(file_path, os.path.join(trash_dir, filename))
    db = SessionLocal()
    try:
        meta = (
            db.query(UserFile)
            .filter_by(username=username, filename=filename)
            .first()
        )
        now = datetime.utcnow()
        if meta:
            meta.deleted_at = now
        else:
            db.add(UserFile(username=username, filename=filename, deleted_at=now))
        db.commit()
    finally:
        db.close()
    delete_share_link(username, filename)
    delete_share_notification(username, filename)
    log_activity(
        username, f"{username} kullanıcısı '{filename}' dosyasını sildi", "delete"
    )
    return jsonify(success=True)


@app.route("/trash/list", methods=["POST"])
def list_trash():
    cleanup_expired_files()
    username = request.form.get("username")
    admin_mode = request.form.get("admin") and is_admin(username)
    db = SessionLocal()
    try:
        query = db.query(UserFile).filter(UserFile.deleted_at != None)
        if not admin_mode:
            query = query.filter_by(username=username)
        metas = query.all()
        now = datetime.utcnow()
        files = []
        for meta in metas:
            remaining = meta.deleted_at + timedelta(days=15) - now
            entry = {
                "filename": meta.filename,
                "time_left": format_remaining(remaining),
            }
            if admin_mode:
                entry["username"] = meta.username
            files.append(entry)
        return jsonify(files=files)
    finally:
        db.close()


@app.route("/trash/empty", methods=["POST"])
def empty_trash():
    username = request.form.get("username")
    admin_mode = request.form.get("admin") and is_admin(username)
    db = SessionLocal()
    try:
        query = db.query(UserFile).filter(UserFile.deleted_at != None)
        if not admin_mode:
            query = query.filter_by(username=username)
        metas = query.all()
        for meta in metas:
            trash_path = os.path.join(DATA_DIR, "_trash", meta.username, meta.filename)
            if os.path.exists(trash_path):
                os.remove(trash_path)
            db.query(ShareLink).filter_by(
                username=meta.username, filename=meta.filename
            ).delete()
            db.delete(meta)
        db.commit()
    finally:
        db.close()
    log_activity(
        username, f"{username} kullanıcısı çöp kutusunu boşalttı", "empty_trash"
    )
    return jsonify(success=True)


@app.route("/trash/restore", methods=["POST"])
def restore_file():
    username = request.form.get("username")
    filename = request.form.get("filename")
    trash_path = os.path.join(DATA_DIR, "_trash", username, filename)
    user_dir = os.path.join(DATA_DIR, username)
    if not os.path.exists(trash_path):
        return jsonify(success=False, error="Dosya bulunamadı")
    os.makedirs(user_dir, exist_ok=True)
    shutil.move(trash_path, os.path.join(user_dir, filename))
    db = SessionLocal()
    try:
        meta = (
            db.query(UserFile)
            .filter_by(username=username, filename=filename)
            .first()
        )
        if meta:
            meta.deleted_at = None
            db.commit()
    finally:
        db.close()
    log_activity(
        username, f"{username} kullanıcısı '{filename}' dosyasını geri aldı", "restore"
    )
    return jsonify(success=True)


@app.route("/file/update", methods=["POST"])
def update_file():
    username = request.form.get("username")
    filename = request.form.get("filename")
    description = request.form.get("description", "")
    share_exp = request.form.get("share_expires_at")
    db = SessionLocal()
    try:
        meta = (
            db.query(UserFile)
            .filter_by(username=username, filename=filename)
            .first()
        )
        if meta:
            meta.description = description
        else:
            db.add(UserFile(username=username, filename=filename, description=description))
        if share_exp is not None:
            link = (
                db.query(ShareLink)
                .filter_by(username=username, filename=filename)
                .first()
            )
            expires_dt = (
                datetime.strptime(share_exp, "%Y-%m-%d") if share_exp else None
            )
            if link:
                link.expires_at = expires_dt
        db.commit()
    finally:
        db.close()
    return jsonify(success=True)


@app.route("/share", methods=["POST"])
def share_file():
    cleanup_expired_files()
    username = request.form.get("username")
    filename = request.form.get("filename")
    user_dir = os.path.join(DATA_DIR, username)
    file_path = os.path.join(user_dir, filename)
    if not os.path.exists(file_path):
        return jsonify(success=False, error="Dosya bulunamadı")
    days = request.form.get("days")
    purpose = request.form.get("purpose", "")
    token = find_share_token(username, filename)
    if token is None:
        token = secrets.token_urlsafe(16)
        expires_at = None
        if days and int(days) > 0:
            expires_at = datetime.utcnow() + timedelta(days=int(days))
        db = SessionLocal()
        try:
            meta = (
                db.query(UserFile)
                .filter_by(username=username, filename=filename)
                .first()
            )
            file_exp = meta.expires_at if meta else None
        finally:
            db.close()
        if file_exp:
            if not expires_at or expires_at > file_exp:
                expires_at = file_exp
        approver_user, approver_email, _ = get_manager_info(username)
        auto_approve = is_admin(username) or not approver_user
        if not approver_user:
            approver_user = next(iter(ADMIN_USERS), None)
            approver_email = get_user_email(approver_user) if approver_user else ""
        approve_token = None
        reject_token = None
        if not auto_approve:
            approve_token = secrets.token_urlsafe(16)
            reject_token = secrets.token_urlsafe(16)
        create_share_link(
            token,
            username,
            filename,
            expires_at,
            approved=auto_approve,
            approve_token=approve_token,
            reject_token=reject_token,
            purpose=purpose,
        )
        log_activity(
            username,
            f"{username} kullanıcısı '{filename}' dosyası için açık paylaşım oluşturdu",
            "share_public",
        )
        if auto_approve:
            create_notification(
                username,
                f"'{filename}' paylaşımı onaylandı",
            )
        else:
            send_approval_email(
                username,
                filename,
                approve_token,
                reject_token,
                approver_email,
                purpose,
            )
            if approver_user:
                create_notification(
                    approver_user,
                    f"'{filename}' dosyası için onay bekleyen paylaşım",
                )
    return jsonify(success=True, link=f"{PUBLIC_BASE_URL}/public/{token}")


@app.route("/share/delete", methods=["POST"])
def delete_share():
    username = request.form.get("username")
    filename = request.form.get("filename")
    delete_share_link(username, filename)
    delete_share_notification(username, filename)
    log_activity(
        username,
        f"{username} kullanıcısı '{filename}' dosyasının paylaşımını kaldırdı",
        "share_public_delete",
    )
    return jsonify(success=True)


@app.route("/share/approve/<token>", methods=["GET"])
def approve_share(token):
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(approve_token=token).first()
        if not link or link.approved or link.rejected:
            return render_template("message.html", message="Geçersiz bağlantı")
        link.approved = True
        link.rejected = False
        link.approve_token = None
        link.reject_token = None
        db.commit()
        create_notification(
            link.username,
            f"'{link.filename}' paylaşımı onaylandı",
        )
        log_activity(
            link.username,
            f"{link.username} kullanıcısının '{link.filename}' paylaşımı bölüm amiri tarafından onaylandı",
            "share_public_approve",
        )
        return render_template("message.html", message="Paylaşım onaylandı")
    finally:
        db.close()


@app.route("/share/reject/<token>", methods=["GET"])
def reject_share(token):
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(reject_token=token).first()
        if not link or link.approved or link.rejected:
            return render_template("message.html", message="Geçersiz bağlantı")
        link.rejected = True
        link.approved = False
        link.approve_token = None
        link.reject_token = None
        db.commit()
        create_notification(
            link.username,
            f"'{link.filename}' paylaşımı reddedildi",
        )
        log_activity(
            link.username,
            f"{link.username} kullanıcısının '{link.filename}' paylaşımı bölüm amiri tarafından reddedildi",
            "share_public_reject",
        )
        return render_template("message.html", message="Paylaşım reddedildi")
    finally:
        db.close()


@app.route("/share/pending", methods=["POST"])
def pending_shares():
    user = request.form.get("username")
    admin_mode = request.form.get("admin") and is_admin(user)
    db = SessionLocal()
    try:
        links = db.query(ShareLink).filter_by(approved=False, rejected=False).all()
        shares = []
        for link in links:
            if admin_mode:
                shares.append(
                    {
                        "token": link.token,
                        "approve_token": link.approve_token,
                        "reject_token": link.reject_token,
                        "username": link.username,
                        "filename": link.filename,
                        "expires_at": link.expires_at.strftime("%d/%m/%Y") if link.expires_at else "",
                        "purpose": link.purpose or "",
                    }
                )
            else:
                mgr_user, _, _ = get_manager_info(link.username)
                if not mgr_user:
                    mgr_user = next(iter(ADMIN_USERS), None)
                if mgr_user == user:
                    shares.append(
                        {
                            "token": link.token,
                            "approve_token": link.approve_token,
                            "reject_token": link.reject_token,
                            "username": link.username,
                            "filename": link.filename,
                            "expires_at": link.expires_at.strftime("%d/%m/%Y") if link.expires_at else "",
                            "purpose": link.purpose or "",
                        }
                    )
        return jsonify(shares=shares)
    finally:
        db.close()


@app.route("/preview/<token>", methods=["GET"])
def preview_file(token):
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(token=token).first()
        if not link:
            return "", 404
        auth_resp = require_manager_auth(link)
        if auth_resp:
            return auth_resp
        file_path = os.path.join(DATA_DIR, link.username, link.filename)
        if not os.path.exists(file_path):
            return "", 404
        if request.args.get("list") == "1":
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(file_path) as zf:
                    return jsonify(files=zf.namelist())
            return jsonify(files=[])
        if request.args.get("download") == "1":
            log_download(link.username, link.filename, session.get("username"))
            return send_file(
                file_path,
                as_attachment=True,
                download_name=link.filename,
            )
        ext = os.path.splitext(file_path)[1].lower()
        office_exts = {".docx", ".xlsx", ".pptx"}
        if ext in office_exts:
            try:
                if ext == ".docx":
                    from docx import Document

                    doc = Document(file_path)
                    html = (
                        "<html><body>"
                        + "".join(f"<p>{p.text}</p>" for p in doc.paragraphs)
                        + "</body></html>"
                    )
                    return Response(html, mimetype="text/html")
                if ext == ".xlsx":
                    from openpyxl import load_workbook

                    wb = load_workbook(file_path, data_only=True)
                    sheet = wb.active
                    html = "<html><body><table border='1'>"
                    for row in sheet.iter_rows(values_only=True):
                        html += "<tr>" + "".join(
                            f"<td>{'' if cell is None else cell}</td>" for cell in row
                        ) + "</tr>"
                    html += "</table></body></html>"
                    return Response(html, mimetype="text/html")
                if ext == ".pptx":
                    from pptx import Presentation

                    prs = Presentation(file_path)
                    parts = ["<html><body>"]
                    for idx, slide in enumerate(prs.slides, start=1):
                        parts.append(f"<h3>Slide {idx}</h3>")
                        for shape in slide.shapes:
                            if getattr(shape, "has_text_frame", False):
                                parts.append(f"<p>{shape.text}</p>")
                    parts.append("</body></html>")
                    return Response("".join(parts), mimetype="text/html")
            except Exception:
                return "", 500
        mime = mimetypes.guess_type(file_path)[0] or "application/octet-stream"
        return send_file(file_path, mimetype=mime)
    finally:
        db.close()


@app.route("/share/user", methods=["POST"])
def share_with_user():
    cleanup_expired_files()
    sender = request.form.get("username")
    recipient = request.form.get("recipient")
    filename = request.form.get("filename")
    expires_at = request.form.get("expires_at")
    expires_dt = datetime.strptime(expires_at, "%Y-%m-%d") if expires_at else None
    user_dir = os.path.join(DATA_DIR, sender)
    file_path = os.path.join(user_dir, filename)
    if not os.path.exists(file_path):
        return jsonify(success=False, error="Dosya bulunamadı")
    db = SessionLocal()
    try:
        meta = (
            db.query(UserFile)
            .filter_by(username=sender, filename=filename)
            .first()
        )
        file_exp = meta.expires_at if meta else None
        if file_exp:
            if not expires_dt or expires_dt > file_exp:
                expires_dt = file_exp
        db.add(
            UserShare(
                sender=sender,
                recipient=recipient,
                filename=filename,
                expires_at=expires_dt,
            )
        )
        db.commit()
        log_activity(
            [sender, recipient],
            f"{sender} kullanıcısı {recipient} kullanıcısına '{filename}' dosyasını paylaştı",
            "share_user",
        )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/incoming", methods=["POST"])
def incoming_files():
    username = request.form.get("username")
    db = SessionLocal()
    try:
        now = datetime.utcnow()
        shares = db.query(UserShare).filter_by(recipient=username).all()
        files = []
        for s in shares:
            if s.expires_at and s.expires_at < now:
                continue
            files.append(
                {
                    "filename": s.filename,
                    "sender": get_full_name(s.sender),
                    "expires_at": s.expires_at.strftime("%d/%m/%Y")
                    if s.expires_at
                    else "",
                }
            )
        return jsonify(files=files)
    finally:
        db.close()


@app.route("/outgoing", methods=["POST"])
def outgoing_files():
    username = request.form.get("username")
    db = SessionLocal()
    try:
        now = datetime.utcnow()
        files = []
        shares = db.query(UserShare).filter_by(sender=username).all()
        for s in shares:
            if s.expires_at and s.expires_at < now:
                continue
            files.append(
                {
                    "filename": s.filename,
                    "target": get_full_name(s.recipient),
                    "target_id": s.recipient,
                    "type": "user",
                    "expires_at": s.expires_at.strftime("%d/%m/%Y")
                    if s.expires_at
                    else "",
                }
            )
        team_shares = db.query(TeamFile).filter_by(username=username).all()
        team_ids = [t.team_id for t in team_shares]
        teams = db.query(Team).filter(Team.id.in_(team_ids)).all() if team_ids else []
        team_names = {t.id: t.name for t in teams}
        for t in team_shares:
            if t.expires_at and t.expires_at < now:
                continue
            files.append(
                {
                    "filename": t.filename,
                    "target": team_names.get(t.team_id, ""),
                    "target_id": t.team_id,
                    "type": "team",
                    "expires_at": t.expires_at.strftime("%d/%m/%Y")
                    if t.expires_at
                    else "",
                }
            )
        return jsonify(files=files)
    finally:
        db.close()


@app.route("/outgoing/delete", methods=["POST"])
def delete_outgoing():
    username = request.form.get("username")
    filename = request.form.get("filename")
    target = request.form.get("target")
    target_type = request.form.get("type")
    db = SessionLocal()
    try:
        if target_type == "user":
            db.query(UserShare).filter_by(
                sender=username, recipient=target, filename=filename
            ).delete()
            db.commit()
            log_activity(
                [username, target],
                f"{username} kullanıcısı {target} kullanıcısına paylaştığı '{filename}' dosyasını kaldırdı",
                "share_user_delete",
            )
        elif target_type == "team":
            team_id = int(target)
            team = db.query(Team).filter_by(id=team_id).first()
            members = (
                db.query(TeamMember)
                .filter_by(team_id=team_id, accepted=True)
                .all()
            )
            db.query(TeamFile).filter_by(
                team_id=team_id, username=username, filename=filename
            ).delete()
            db.commit()
            member_usernames = [m.username for m in members]
            if username not in member_usernames:
                member_usernames.append(username)
            team_name = team.name if team else ""
            log_activity(
                member_usernames,
                f"{username} kullanıcısı {team_name} ekibinden '{filename}' dosyasını sildi",
                "share_team_delete",
            )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/public/<token>", methods=["GET"])
def public_page(token):
    cleanup_expired_files()
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(token=token).first()
        if link and link.expires_at and link.expires_at < datetime.utcnow():
            db.delete(link)
            db.commit()
            link = None
    finally:
        db.close()
    if not link:
        return render_template("public.html", error="Bağlantının süresi dolmuş.")
    if link.rejected:
        return render_template("public.html", error="Bağlantı reddedildi.")
    if not link.approved:
        return render_template("public.html", error="Bağlantı onay bekliyor.")
    username = link.username
    filename = link.filename
    user_dir = os.path.join(DATA_DIR, username)
    file_path = os.path.join(user_dir, filename)
    if not os.path.exists(file_path):
        return render_template("public.html", error="Dosya sunucudan kaldırılmıştır.")
    size = format_file_size(os.path.getsize(file_path))
    uploader = get_full_name(username)
    db = SessionLocal()
    try:
        meta = db.query(UserFile).filter_by(username=username, filename=filename).first()
        description = meta.description if meta else ""
    finally:
        db.close()
    return render_template(
        "public.html",
        token=token,
        filename=filename,
        uploader=uploader,
        size=size,
        description=description,
        expires_at=link.expires_at.strftime("%d/%m/%Y") if link.expires_at else "",
    )


@app.route("/public/<token>/download", methods=["GET"])
def public_download_file(token):
    cleanup_expired_files()
    db = SessionLocal()
    try:
        link = db.query(ShareLink).filter_by(token=token).first()
        if link and link.expires_at and link.expires_at < datetime.utcnow():
            db.delete(link)
            db.commit()
            link = None
    finally:
        db.close()
    if not link:
        return render_template("public.html", error="Bağlantının süresi dolmuş.")
    if link.rejected:
        return render_template("public.html", error="Bağlantı reddedildi.")
    if not link.approved:
        return render_template("public.html", error="Bağlantı onay bekliyor.")
    username = link.username
    filename = link.filename
    user_dir = os.path.join(DATA_DIR, username)
    file_path = os.path.join(user_dir, filename)
    if not os.path.exists(file_path):
        return render_template("public.html", error="Dosya sunucudan kaldırılmıştır.")
    log_download(username, filename, session.get("username"))
    return send_file(file_path, as_attachment=True, download_name=filename)


@app.route("/stats", methods=["POST"])
def stats():
    username = request.form.get("username")
    user_dir = os.path.join(DATA_DIR, username)
    file_count = len(os.listdir(user_dir)) if os.path.exists(user_dir) else 0
    db = SessionLocal()
    try:
        logs = db.query(DownloadLog).filter_by(username=username).all()
        logs_data = [
            {
                "filename": l.filename,
                "timestamp": l.timestamp.isoformat(),
                "ip_address": l.ip_address,
                "country": l.country,
            }
            for l in logs
        ]
        counts = {}
        for l in logs:
            counts[l.filename] = counts.get(l.filename, 0) + 1
    finally:
        db.close()
    return jsonify(
        file_count=file_count,
        download_count=len(logs_data),
        download_logs=logs_data,
        download_counts=counts,
    )


@app.route("/dashboard/data", methods=["POST"])
def dashboard_data():
    username = request.form.get("username")
    user_dir = os.path.join(DATA_DIR, username)
    files = []
    total_size = 0
    if os.path.exists(user_dir):
        for name in os.listdir(user_dir):
            path = os.path.join(user_dir, name)
            if os.path.isfile(path):
                size = os.path.getsize(path)
                files.append({"name": name, "size": size})
                total_size += size
    files.sort(key=lambda f: f["size"], reverse=True)
    top_disk_files = files[:5]
    file_count = len(files)
    db = SessionLocal()
    try:
        logs = db.query(DownloadLog).filter_by(username=username).all()
        counts_by_file = {}
        counts_by_country = {}
        for l in logs:
            counts_by_file[l.filename] = counts_by_file.get(l.filename, 0) + 1
            if l.country:
                counts_by_country[l.country] = counts_by_country.get(l.country, 0) + 1
        top_files = sorted(
            counts_by_file.items(), key=lambda x: x[1], reverse=True
        )[:5]
        top_countries = sorted(
            counts_by_country.items(), key=lambda x: x[1], reverse=True
        )[:5]
        download_count = len(logs)

        memberships = (
            db.query(TeamMember)
            .filter_by(username=username, accepted=True)
            .all()
        )
        team_ids = [m.team_id for m in memberships]
        teams = []
        if team_ids:
            for team in db.query(Team).filter(Team.id.in_(team_ids)).all():
                members = (
                    db.query(TeamMember)
                    .filter_by(team_id=team.id, accepted=True)
                    .all()
                )
                file_count = (
                    db.query(TeamFile).filter_by(team_id=team.id).count()
                )
                teams.append(
                    {
                        "id": team.id,
                        "name": team.name,
                        "member_count": len(members),
                        "members": [get_full_name(m.username) for m in members],
                        "file_count": file_count,
                    }
                )

        now = datetime.utcnow()
        upcoming = now + timedelta(days=7)
        user_files = db.query(UserFile).filter_by(username=username).all()
        user_links = db.query(ShareLink).filter_by(username=username).all()
        expiring_map = {}
        for f in user_files:
            if f.expires_at and now < f.expires_at <= upcoming:
                delta = f.expires_at - now
                expiring_map[f.filename] = min(
                    expiring_map.get(f.filename, delta), delta
                )
        for l in user_links:
            if (
                l.expires_at
                and now < l.expires_at <= upcoming
                and l.approved
                and not l.rejected
            ):
                delta = l.expires_at - now
                expiring_map[l.filename] = min(
                    expiring_map.get(l.filename, delta), delta
                )
        expiring = [(fn, td) for fn, td in expiring_map.items()]
        expiring.sort(key=lambda x: x[1])
        expiring = [
            {"filename": fn, "time_left": format_remaining(td)} for fn, td in expiring
        ]
    finally:
        db.close()

    return jsonify(
        disk_usage={"total": total_size, "files": top_disk_files},
        file_count=file_count,
        download_count=download_count,
        downloads={
            "files": [{"filename": f, "count": c} for f, c in top_files],
            "countries": [{"country": c, "count": n} for c, n in top_countries],
        },
        teams=teams,
        expiring_files=expiring,
    )


@app.route("/teams/create", methods=["POST"])
def create_team():
    username = request.form.get("username")
    team_name = request.form.get("team_name")
    members = request.form.get("members", "")
    member_list = [m.strip() for m in members.split(",") if m.strip()]
    db = SessionLocal()
    try:
        team = Team(name=team_name, creator=username)
        db.add(team)
        db.commit()
        db.refresh(team)
        db.add(TeamMember(team_id=team.id, username=username, accepted=True))
        for m in member_list:
            db.add(TeamMember(team_id=team.id, username=m, accepted=False))
            create_notification(
                m,
                f"{username} adlı kullanıcı seni {team_name} ekibine davet etti.",
                team_id=team.id,
            )
        db.commit()
        log_activity(
            username,
            f"{username} kullanıcısı {team_name} ekibini oluşturdu",
            "team_create",
        )
        return jsonify(success=True, team_id=team.id)
    finally:
        db.close()


@app.route("/teams/delete", methods=["POST"])
def delete_team():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    db = SessionLocal()
    try:
        team = db.query(Team).filter_by(id=team_id).first()
        if not team:
            return jsonify(success=False, error="Ekip bulunamadı")
        if team.creator != username and not is_admin(username):
            return jsonify(success=False, error="Yetkiniz yok")
        members = (
            db.query(TeamMember).filter_by(team_id=team_id, accepted=True).all()
        )
        member_usernames = [m.username for m in members]
        if username not in member_usernames:
            member_usernames.append(username)
        team_name = team.name
        db.query(TeamMember).filter_by(team_id=team_id).delete()
        db.query(TeamFile).filter_by(team_id=team_id).delete()
        db.query(Notification).filter_by(team_id=team_id).delete()
        db.delete(team)
        db.commit()
        log_activity(
            member_usernames,
            f"{username} kullanıcısı {team_name} ekibini sildi",
            "team_delete",
        )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/leave", methods=["POST"])
def leave_team():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    db = SessionLocal()
    try:
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=username)
            .first()
        )
        if membership:
            team = db.query(Team).filter_by(id=team_id).first()
            members = (
                db.query(TeamMember)
                .filter_by(team_id=team_id, accepted=True)
                .all()
            )
            member_usernames = [m.username for m in members]
            if username not in member_usernames:
                member_usernames.append(username)
            team_name = team.name if team else ""
            db.delete(membership)
            db.commit()
            log_activity(
                member_usernames,
                f"{username} kullanıcısı {team_name} ekibinden ayrıldı",
                "team_leave",
            )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/remove_member", methods=["POST"])
def remove_member():
    requester = request.form.get("username")
    team_id = request.form.get("team_id")
    member = request.form.get("member")
    db = SessionLocal()
    try:
        team = db.query(Team).filter_by(id=team_id).first()
        if not team:
            return jsonify(success=False, error="Ekip bulunamadı")
        if member == team.creator:
            return jsonify(success=False, error="Kurucu silinemez")
        if requester != member and team.creator != requester:
            return jsonify(success=False, error="Yetkiniz yok")
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=member)
            .first()
        )
        if membership:
            members = (
                db.query(TeamMember)
                .filter_by(team_id=team_id, accepted=True)
                .all()
            )
            member_usernames = [m.username for m in members]
            for u in [member, requester]:
                if u not in member_usernames:
                    member_usernames.append(u)
            team_name = team.name
            db.delete(membership)
            db.commit()
            log_activity(
                member_usernames,
                f"{requester} kullanıcısı {team_name} ekibinden {member} kullanıcısını çıkardı",
                "team_remove_member",
            )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/list", methods=["POST"])
def list_teams():
    username = request.form.get("username")
    admin_mode = request.form.get("admin") and is_admin(username)
    db = SessionLocal()
    try:
        if admin_mode:
            teams = db.query(Team).all()
        else:
            memberships = db.query(TeamMember).filter_by(username=username).all()
            team_ids = [m.team_id for m in memberships]
            teams = db.query(Team).filter(Team.id.in_(team_ids)).all()
        data = []
        for t in teams:
            members = db.query(TeamMember).filter_by(team_id=t.id).all()
            data.append(
                {
                    "id": t.id,
                    "name": t.name,
                    "creator": t.creator,
                    "members": [
                        {
                            "username": m.username,
                            "name": get_full_name(m.username),
                            "accepted": m.accepted,
                        }
                        for m in members
                    ],
                }
            )
        return jsonify(teams=data)
    finally:
        db.close()


@app.route("/teams/add_files", methods=["POST"])
def add_files_to_team():
    cleanup_expired_files()
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    filenames = request.form.getlist("filenames")
    expires_at = request.form.get("expires_at")
    expires_dt = datetime.strptime(expires_at, "%Y-%m-%d") if expires_at else None
    db = SessionLocal()
    try:
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=username)
            .first()
        )
        if not membership and not is_admin(username):
            return jsonify(success=False, error="Yetkiniz yok")
        team = db.query(Team).filter_by(id=team_id).first()
        members = db.query(TeamMember).filter_by(team_id=team_id, accepted=True).all()
        for fname in filenames:
            file_path = os.path.join(DATA_DIR, username, fname)
            if not os.path.exists(file_path):
                continue
            meta = (
                db.query(UserFile)
                .filter_by(username=username, filename=fname)
                .first()
            )
            file_exp = meta.expires_at if meta else None
            share_exp = expires_dt
            if file_exp:
                if not share_exp or share_exp > file_exp:
                    share_exp = file_exp
            db.add(
                TeamFile(
                    team_id=team_id,
                    username=username,
                    filename=fname,
                    expires_at=share_exp,
                )
            )
        db.commit()
        team_name = team.name if team else ""
        member_usernames = [m.username for m in members]
        for fname in filenames:
            log_activity(
                member_usernames,
                f"{username} kullanıcısı {team_name} ekibine '{fname}' dosyasını yükledi",
                "team_add_file",
            )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/delete_file", methods=["POST"])
def delete_team_file():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    filename = request.form.get("filename")
    db = SessionLocal()
    try:
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=username)
            .first()
        )
        team = db.query(Team).filter_by(id=team_id).first()
        if not membership and not is_admin(username):
            return jsonify(success=False, error="Yetkiniz yok")
        db.query(TeamFile).filter_by(team_id=team_id, filename=filename).delete()
        db.commit()
        members = db.query(TeamMember).filter_by(team_id=team_id, accepted=True).all()
        member_usernames = [m.username for m in members]
        team_name = team.name if team else ""
        log_activity(
            member_usernames,
            f"{username} kullanıcısı {team_name} ekibinden '{filename}' dosyasını sildi",
            "team_delete_file",
        )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/details", methods=["POST"])
def team_details():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    db = SessionLocal()
    try:
        team = db.query(Team).filter_by(id=team_id).first()
        if not team:
            return jsonify(success=False, error="Ekip bulunamadı")
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=username)
            .first()
        )
        if not membership and not is_admin(username):
            return jsonify(success=False, error="Yetkiniz yok")
        members = db.query(TeamMember).filter_by(team_id=team_id).all()
        files = db.query(TeamFile).filter_by(team_id=team_id).all()
        return jsonify(
            success=True,
            team={
                "id": team.id,
                "name": team.name,
                "creator": team.creator,
                "members": [
                    {
                        "username": m.username,
                        "name": get_full_name(m.username),
                        "accepted": m.accepted,
                    }
                    for m in members
                ],
                "files": [
                    {"filename": f.filename, "username": f.username} for f in files
                ],
            },
        )
    finally:
        db.close()


@app.route("/teams/add_member", methods=["POST"])
def add_member_to_team():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    new_member = request.form.get("new_member")
    db = SessionLocal()
    try:
        team = db.query(Team).filter_by(id=team_id).first()
        if not team:
            return jsonify(success=False, error="Ekip bulunamadı")
        if team.creator != username:
            return jsonify(success=False, error="Yetkiniz yok")
        exists = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=new_member)
            .first()
        )
        if exists:
            return jsonify(success=False, error="Kullanıcı zaten ekipte")
        db.add(TeamMember(team_id=team_id, username=new_member, accepted=False))
        db.commit()
        create_notification(
            new_member,
            f"{username} adlı kullanıcı seni {team.name} ekibine davet etti.",
            team_id=team.id,
        )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/accept", methods=["POST"])
def accept_team():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    db = SessionLocal()
    try:
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=username)
            .first()
        )
        if not membership:
            return jsonify(success=False, error="Ekip bulunamadı")
        membership.accepted = True
        db.commit()
        team = db.query(Team).filter_by(id=team_id).first()
        members = db.query(TeamMember).filter_by(team_id=team_id, accepted=True).all()
        member_usernames = [m.username for m in members]
        team_name = team.name if team else ""
        if team:
            create_notification(
                team.creator,
                f"{username} kullanıcısı {team.name} ekibine katılma davetini kabul etti.",
            )
        log_activity(
            member_usernames,
            f"{username} kullanıcısı {team_name} ekibine katıldı",
            "team_accept",
        )
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/teams/reject", methods=["POST"])
def reject_team():
    username = request.form.get("username")
    team_id = request.form.get("team_id")
    db = SessionLocal()
    try:
        membership = (
            db.query(TeamMember)
            .filter_by(team_id=team_id, username=username)
            .first()
        )
        team = db.query(Team).filter_by(id=team_id).first()
        if membership:
            db.delete(membership)
            db.commit()
            if team:
                msg = f"{username} kullanıcısı {team.name} ekibine katılma davetini reddetti"
                create_notification(team.creator, msg)
                log_activity([team.creator, username], msg, "team_reject")
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/notifications/clear", methods=["POST"])
def clear_notifications():
    username = request.form.get("username")
    db = SessionLocal()
    try:
        db.query(Notification).filter_by(username=username).delete()
        db.commit()
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/notifications", methods=["POST"])
def notifications():
    username = request.form.get("username")
    db = SessionLocal()
    try:
        notifs = (
            db.query(Notification)
            .filter_by(username=username)
            .order_by(Notification.created_at.desc())
            .all()
        )
        data = [
            {
                "id": n.id,
                "message": n.message,
                "created_at": n.created_at.strftime("%d/%m/%Y %H:%M"),
                "read": n.read,
                "team_id": n.team_id,
            }
            for n in notifs
        ]
        return jsonify(notifications=data)
    finally:
        db.close()


@app.route("/notifications/read", methods=["POST"])
def read_notifications():
    username = request.form.get("username")
    ids = request.form.getlist("ids[]") or request.form.getlist("ids")
    db = SessionLocal()
    try:
        if ids:
            q = db.query(Notification).filter(
                Notification.username == username, Notification.id.in_(ids)
            )
            q.update({"read": True}, synchronize_session=False)
            db.commit()
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/notifications/delete", methods=["POST"])
def delete_notifications():
    username = request.form.get("username")
    ids = request.form.getlist("ids[]") or request.form.getlist("ids")
    db = SessionLocal()
    try:
        if ids:
            db.query(Notification).filter(
                Notification.username == username, Notification.id.in_(ids)
            ).delete(synchronize_session=False)
            db.commit()
        return jsonify(success=True)
    finally:
        db.close()


@app.route("/activities", methods=["POST"])
def activities():
    username = request.form.get("username")
    categories = request.form.get("categories")
    users = request.form.get("users")
    db = SessionLocal()
    try:
        base_query = db.query(Activity)
        if not is_admin(username):
            base_query = base_query.filter_by(username=username)

        all_acts = base_query.all()
        cat_set = sorted({a.category for a in all_acts if a.category})
        user_set = sorted({a.username for a in all_acts})

        if categories:
            cat_list = [c for c in categories.split(",") if c]
            base_query = base_query.filter(Activity.category.in_(cat_list))
        if users and is_admin(username):
            user_list = [u for u in users.split(",") if u]
            base_query = base_query.filter(Activity.username.in_(user_list))

        acts = base_query.order_by(Activity.created_at.desc()).all()
        data = []
        for a in acts:
            msg = a.message
            prefix = f"{a.username} kullanıcısı "
            if msg.startswith(prefix):
                msg = msg[len(prefix):]
            data.append(
                {
                    "username": a.username,
                    "message": msg,
                    "category": a.category,
                    "created_at": a.created_at.strftime("%d/%m/%Y %H:%M"),
                }
            )
        return jsonify(activities=data, categories=cat_set, users=user_set)
    finally:
        db.close()



if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8000)

